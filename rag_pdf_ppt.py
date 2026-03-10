from PyPDF2 import PdfReader
from pptx import Presentation

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_ollama import OllamaLLM


# -------------------------------
# 1. Load Document (PDF or PPT)
# -------------------------------
def load_document(file_path):
    text = ""

    # If PDF
    if file_path.lower().endswith(".pdf"):
        reader = PdfReader(file_path)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted

    # If PPT or PPTX
    elif file_path.lower().endswith((".ppt", ".pptx")):
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    else:
        raise ValueError("Unsupported file format. Use PDF or PPT.")

    return text


# -------------------------------
# 2. Split text into chunks
# -------------------------------
def split_text(text):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=100
    )
    return splitter.split_text(text)


# -------------------------------
# 3. Create Vector Store
# -------------------------------
def create_vector_store(chunks):
    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )

    vector_db = FAISS.from_texts(chunks, embeddings)
    return vector_db


# -------------------------------
# 4. Answer Question (Manual RAG)
# -------------------------------
def answer_question(vector_db, query):
    docs = vector_db.similarity_search(query, k=3)
    context = "\n\n".join([doc.page_content for doc in docs])

    llm = OllamaLLM(model="llama3")

    prompt = f"""
You are an AI assistant answering questions based ONLY on the given document context.

Context:
{context}

Question:
{query}

Answer clearly and accurately:
"""

    return llm.invoke(prompt)


# -------------------------------
# 5. Main Chat Loop
# -------------------------------
def chat_with_document(file_path):
    print("📄 Loading document...")
    text = load_document(file_path)

    print("✂️ Splitting text...")
    chunks = split_text(text)

    print("🧠 Creating vector database...")
    vector_db = create_vector_store(chunks)

    print("\n✅ Document ready! Ask questions (type 'exit' to stop)\n")

    while True:
        query = input("You: ")

        if query.lower() == "exit":
            print("👋 Exiting chat...")
            break

        response = answer_question(vector_db, query)
        print("\nAI:", response, "\n")


# -------------------------------
# Run
# -------------------------------
if __name__ == "__main__":
    file_path = r"C:\pc\PROJECT_DRIVE\PROJECTS\LLM_RAG\SCANSTORM.pdf"
    # You can also give:
    # r"C:\pc\PROJECT_DRIVE\PROJECTS\LLM_RAG\presentation.pptx"

    chat_with_document(file_path)